import docx
import fitz  # PyMuPDF
import pandas as pd
import pptx
import os
from .data_preprocessing import clean_text

def extract_data_from_docx(file_path):
    """Extract text from a DOCX file."""
    doc = docx.Document(file_path)
    text = []
    for para in doc.paragraphs:
        text.append(para.text)
    return "\n".join(text)

def extract_data_from_pdf(file_path):
    """Extract text from a PDF file using PyMuPDF."""
    doc = fitz.open(file_path)
    text = []
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        text.append(page.get_text())
    return "\n".join(text)

def extract_data_from_xlsx(file_path):
    """Extract text from an XLSX file using pandas."""
    # Use pandas to read the Excel file
    df = pd.read_excel(file_path, sheet_name=None)  # Read all sheets
    text = []
    for sheet_name, sheet_data in df.items():
        text.append(f"Sheet: {sheet_name}")
        text.append(sheet_data.to_string(index=False))  # Convert the sheet to text
    return "\n".join(text)

def extract_data_from_pptx(file_path):
    """Extract text from a PPTX file."""
    ppt = pptx.Presentation(file_path)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_data(file_path, file_extension):
    """Choose extraction function based on file type."""
    file_extension = file_extension.lstrip(".").lower()  # Normalize extension

    if file_extension == "docx":
        return extract_data_from_docx(file_path)
    elif file_extension == "pdf":
        return extract_data_from_pdf(file_path)
    elif file_extension == "xlsx":
        return extract_data_from_xlsx(file_path)
    elif file_extension == "pptx":
        return extract_data_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")


def extract_and_preprocess(file_path, file_extension):
    """Extract and preprocess text from different file formats."""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    extracted_text = extract_data(file_path, file_extension)
    cleaned_text = clean_text(extracted_text)
    # print(cleaned_text)
    return cleaned_text